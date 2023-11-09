# process_pptx.py
from pptx import Presentation
import os
from utils import setup_directories
import logging

def extract_text_and_images_pptx(pptx_path):
    try:
        base_dir = "output"
        pptx_dir_name = os.path.splitext(os.path.basename(pptx_path))[0]
        pptx_base_dir = os.path.join(base_dir, pptx_dir_name)
        directories = setup_directories(pptx_base_dir)

        prs = Presentation(pptx_path)
        text_content = []

        for slide_number, slide in enumerate(prs.slides, start=1):
            # Extract text from each slide
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)

            # Save images from each slide
            for shape in slide.shapes:
                if shape.shape_type == 13:  # This is the type number for Picture
                    image = shape.image
                    image_bytes = image.blob
                    image_filename = f"image_slide_{slide_number}.jpeg"
                    image_path = os.path.join(directories["Images"], image_filename)
                    with open(image_path, 'wb') as f:
                        f.write(image_bytes)

        text_filename = os.path.join(directories["Texts"], "extracted_text.txt")
        with open(text_filename, 'w', encoding="utf-8") as f:
            f.write("\n\n".join(text_content))

        logging.info(f"Processed {pptx_path} successfully.")

    except Exception as e:
        logging.error(f"Error processing {pptx_path}: {e}")
